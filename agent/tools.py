import os
import re
import csv
import json
import logging
import httpx
import trafilatura
from bs4 import BeautifulSoup
import discord
from typing import Any
from tools.registry import tool_registry, ToolExecutionContext
from core.searxng_client import searxng_client
from core.client_manager import client_manager
from agent.session_manager import session_manager
from google.genai import types

logger = logging.getLogger("PriestyAI.Agent.Tools")

MAX_TOOL_OUTPUT_CHARS = 4500

def is_youtube_url(url: str) -> bool:
    return bool(re.search(r'(?:youtube\.com\/(?:watch\?v=|shorts\/|embed\/)|youtu\.be\/)', url, re.IGNORECASE))

def normalize_youtube_url(url: str) -> str:
    match = re.search(r'youtu\.be\/([a-zA-Z0-9_-]+)', url)
    if match:
        return f"https://www.youtube.com/watch?v={match.group(1)}"
    return url

def _resolve_safe_workspace_path(workspace_root: str, relative_path: str) -> str | None:
    norm_root = os.path.abspath(workspace_root)
    clean_rel = relative_path.strip().lstrip("./").lstrip("/")
    target = os.path.abspath(os.path.join(norm_root, clean_rel))
    try:
        if os.path.commonpath([norm_root, target]) == norm_root:
            return target
    except ValueError:
        return None
    return None

@tool_registry.register(
    name="agent_terminal",
    description=(
        "Executes a shell command inside the session's persistent workspace container.\n"
        "- command: The shell command to run (e.g. 'npm test', 'pytest', 'pip install -r requirements.txt', 'cargo build')."
    )
)
async def agent_terminal(command: str, context: ToolExecutionContext = None) -> dict[str, Any]:
    if not context or not getattr(context, "agent_session_id", None):
        return {"error": "Active agent session context is required."}

    session_id = context.agent_session_id
    code, stdout, stderr = await session_manager.exec_in_container(session_id, command)
    return {
        "command": command,
        "exit_code": code,
        "stdout": (stdout or "(no output)")[:MAX_TOOL_OUTPUT_CHARS],
        "stderr": (stderr or "")[:1500] if stderr else None
    }

@tool_registry.register(
    name="agent_read_file",
    description=(
        "Reads lines and structured contents from files in the workspace.\n"
        "Supports code, markdown, PDFs (with tables), Word (.docx), Excel spreadsheets (.xlsx/.xls), "
        "CSV/TSV data tables, PowerPoint slides (.pptx), and Jupyter Notebooks (.ipynb).\n"
        "- path: Relative path from workspace root.\n"
        "- start_line: 1-indexed starting line (default 1).\n"
        "- end_line: Ending line (0 for full content)."
    )
)
async def agent_read_file(path: str, start_line: int = 1, end_line: int = 0, context: ToolExecutionContext = None) -> dict[str, Any]:
    if not context or not getattr(context, "agent_session_id", None):
        return {"error": "Active agent session context is required."}

    session = session_manager.get_session_by_id(context.agent_session_id)
    if not session:
        return {"error": "Session record missing."}

    full_path = _resolve_safe_workspace_path(session["workspace_path"], path)
    if not full_path:
        return {"error": "Access denied: path outside workspace."}

    if not os.path.exists(full_path):
        return {"error": f"File '{path}' does not exist in workspace."}

    try:
        lower_path = path.lower()

        if lower_path.endswith((".xlsx", ".xls")):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(full_path, data_only=True)
                sheet_blocks = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_blocks.append(f"### Sheet: {sheet_name}")
                    rows = list(sheet.iter_rows(values_only=True))
                    non_empty_rows = [r for r in rows if any(cell is not None for cell in r)]
                    if non_empty_rows:
                        headers = [str(c or "").strip().replace("\n", " ") for c in non_empty_rows[0]]
                        sheet_blocks.append("| " + " | ".join(headers) + " |")
                        sheet_blocks.append("| " + " | ".join(["---"] * len(headers)) + " |")
                        for r in non_empty_rows[1:100]:
                            row_vals = [str(c or "").strip().replace("\n", " ") for c in r]
                            sheet_blocks.append("| " + " | ".join(row_vals) + " |")
                        if len(non_empty_rows) > 100:
                            sheet_blocks.append(f"-# ... and {len(non_empty_rows) - 100} more rows")
                    else:
                        sheet_blocks.append("*(Empty sheet)*")
                
                full_text = "\n\n".join(sheet_blocks)
                lines = [l + "\n" for l in full_text.splitlines()]
            except Exception as ex:
                return {"error": f"Failed to read Excel workbook: {ex}"}

        elif lower_path.endswith((".csv", ".tsv")):
            delimiter = "\t" if lower_path.endswith(".tsv") else ","
            with open(full_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f, delimiter=delimiter)
                rows = list(reader)

            table_blocks = []
            if rows:
                headers = [c.strip().replace("\n", " ") for c in rows[0]]
                table_blocks.append("| " + " | ".join(headers) + " |")
                table_blocks.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for r in rows[1:120]:
                    row_vals = [c.strip().replace("\n", " ") for c in r]
                    table_blocks.append("| " + " | ".join(row_vals) + " |")
                if len(rows) > 120:
                    table_blocks.append(f"-# ... and {len(rows) - 120} more rows")
            
            full_text = "\n".join(table_blocks)
            lines = [l + "\n" for l in full_text.splitlines()]

        elif lower_path.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(full_path)
                slide_blocks = []
                for idx, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            texts.append(shape.text.strip())
                    
                    slide_str = f"--- [Slide {idx + 1}] ---\n" + "\n".join(texts)
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                        slide_str += f"\n**Speaker Notes:** {slide.notes_slide.notes_text_frame.text.strip()}"
                    slide_blocks.append(slide_str)

                full_text = "\n\n".join(slide_blocks)
                lines = [l + "\n" for l in full_text.splitlines()]
            except Exception as ex:
                return {"error": f"Failed to read PowerPoint file: {ex}"}

        elif lower_path.endswith(".ipynb"):
            with open(full_path, "r", encoding="utf-8", errors="replace") as f:
                nb_data = json.load(f)
            
            cells = nb_data.get("cells", [])
            cell_blocks = []
            for idx, c in enumerate(cells):
                c_type = c.get("cell_type", "code")
                src = "".join(c.get("source", []))
                if c_type == "markdown":
                    cell_blocks.append(f"<!-- Markdown Cell {idx + 1} -->\n{src}")
                else:
                    cell_blocks.append(f"```python\n# [Code Cell {idx + 1}]\n{src}\n```")
                    outputs = c.get("outputs", [])
                    out_texts = []
                    for out in outputs:
                        if "text" in out:
                            out_texts.append("".join(out["text"]))
                        elif "data" in out and "text/plain" in out["data"]:
                            out_texts.append("".join(out["data"]["text/plain"]))
                    if out_texts:
                        cell_blocks.append("```text\n# [Output]\n" + "\n".join(out_texts)[:800] + "\n```")

            full_text = "\n\n".join(cell_blocks)
            lines = [l + "\n" for l in full_text.splitlines()]

        elif lower_path.endswith(".pdf"):
            try:
                import pdfplumber
                page_texts = []
                with pdfplumber.open(full_path) as pdf:
                    for idx, page in enumerate(pdf.pages):
                        page_blocks = [f"--- [PDF Page {idx + 1}] ---"]
                        tables = page.extract_tables()
                        if tables:
                            page_blocks.append("\n**[Extracted Tables]**")
                            for table in tables:
                                clean_table = [[(cell or "").strip().replace("\n", " ") for cell in row] for row in table if any(row)]
                                if clean_table:
                                    headers = clean_table[0]
                                    page_blocks.append("\n| " + " | ".join(headers) + " |")
                                    page_blocks.append("| " + " | ".join(["---"] * len(headers)) + " |")
                                    for row in clean_table[1:]:
                                        page_blocks.append("| " + " | ".join(row) + " |")
                                    page_blocks.append("")

                        raw_text = page.extract_text() or ""
                        if raw_text:
                            page_blocks.append("\n**[Page Text]**\n" + raw_text)
                            
                        page_texts.append("\n".join(page_blocks))
                
                full_text = "\n\n".join(page_texts)
                lines = [l + "\n" for l in full_text.splitlines()]
            except ImportError:
                import pypdf
                reader = pypdf.PdfReader(full_path)
                page_texts = []
                for idx, page in enumerate(reader.pages):
                    ptxt = page.extract_text() or ""
                    page_texts.append(f"--- [PDF Page {idx + 1}] ---\n{ptxt}")
                full_text = "\n\n".join(page_texts)
                lines = [l + "\n" for l in full_text.splitlines()]

        elif lower_path.endswith(".docx"):
            try:
                import docx
                doc = docx.Document(full_path)
                doc_lines = []
                for elem in doc.element.body:
                    if elem.tag.endswith('p'):
                        p = docx.text.paragraph.Paragraph(elem, doc)
                        if p.text.strip():
                            doc_lines.append(p.text)
                    elif elem.tag.endswith('tbl'):
                        table = docx.table.Table(elem, doc)
                        table_data = []
                        for row in table.rows:
                            row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                            table_data.append(row_cells)
                        if table_data:
                            doc_lines.append("\n| " + " | ".join(table_data[0]) + " |")
                            doc_lines.append("| " + " | ".join(["---"] * len(table_data[0])) + " |")
                            for row in table_data[1:]:
                                doc_lines.append("| " + " | ".join(row) + " |")
                            doc_lines.append("")
                
                lines = [l + "\n" for l in doc_lines]
            except Exception:
                with open(full_path, "r", encoding="utf-8", errors="replace") as f:
                    lines = f.readlines()

        else:
            with open(full_path, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()

        total = len(lines)
        s_line = max(1, start_line)
        e_line = min(total, end_line) if end_line > 0 else total
        sliced = lines[s_line - 1:e_line]

        formatted = "".join([f"{s_line + idx:4d} | {l}" for idx, l in enumerate(sliced)])
        return {
            "path": path,
            "total_lines": total,
            "showing_lines": f"{s_line}-{e_line}",
            "content": formatted[:MAX_TOOL_OUTPUT_CHARS]
        }
    except Exception as e:
        return {"error": f"Failed to read file '{path}': {e}"}

@tool_registry.register(
    name="agent_write_file",
    description=(
        "Creates a new file or overwrites an entire file in the workspace with complete content.\n"
        "- path: Relative path from workspace root.\n"
        "- content: Full source content to write."
    )
)
async def agent_write_file(path: str, content: str, context: ToolExecutionContext = None) -> dict[str, Any]:
    if not context or not getattr(context, "agent_session_id", None):
        return {"error": "Active agent session context is required."}

    session = session_manager.get_session_by_id(context.agent_session_id)
    if not session:
        return {"error": "Session record missing."}

    full_path = _resolve_safe_workspace_path(session["workspace_path"], path)
    if not full_path:
        return {"error": "Access denied: path outside workspace."}

    os.makedirs(os.path.dirname(full_path), exist_ok=True)
    try:
        with open(full_path, "w", encoding="utf-8") as f:
            f.write(content)
        lines = len(content.splitlines())
        return {"status": "written", "path": path, "lines": lines}
    except Exception as e:
        return {"error": f"Failed to write file: {e}"}

@tool_registry.register(
    name="agent_edit_diff",
    description=(
        "Performs a surgical search-and-replace edit on an existing workspace file.\n"
        "- path: Relative path from workspace root.\n"
        "- search_block: Exact block of lines to find.\n"
        "- replace_block: Replacement block of lines."
    )
)
async def agent_edit_diff(path: str, search_block: str, replace_block: str, context: ToolExecutionContext = None) -> dict[str, Any]:
    if not context or not getattr(context, "agent_session_id", None):
        return {"error": "Active agent session context is required."}

    session = session_manager.get_session_by_id(context.agent_session_id)
    if not session:
        return {"error": "Session record missing."}

    full_path = _resolve_safe_workspace_path(session["workspace_path"], path)
    if not full_path:
        return {"error": "Access denied: path outside workspace."}

    if not os.path.exists(full_path):
        return {"error": f"File '{path}' does not exist."}

    try:
        with open(full_path, "r", encoding="utf-8") as f:
            original = f.read()

        if search_block in original:
            new_content = original.replace(search_block, replace_block, 1)
        else:
            clean_search = search_block.strip()
            if clean_search and clean_search in original:
                new_content = original.replace(clean_search, replace_block.strip(), 1)
            else:
                return {"error": "search_block not found in file. Ensure exact line and code matches."}

        with open(full_path, "w", encoding="utf-8") as f:
            f.write(new_content)

        return {"status": "patched", "path": path}
    except Exception as e:
        return {"error": f"Failed to apply diff: {e}"}

@tool_registry.register(
    name="agent_list_dir",
    description="Lists files and subdirectories in the workspace (subpath relative to root, e.g. '' or 'src')."
)
async def agent_list_dir(subpath: str = "", context: ToolExecutionContext = None) -> dict[str, Any]:
    if not context or not getattr(context, "agent_session_id", None):
        return {"error": "Active agent session context is required."}

    session = session_manager.get_session_by_id(context.agent_session_id)
    if not session:
        return {"error": "Session record missing."}

    workspace_root = session["workspace_path"]
    target_dir = _resolve_safe_workspace_path(workspace_root, subpath or ".")
    if not target_dir or not os.path.exists(target_dir):
        return {"error": f"Directory '{subpath}' does not exist in workspace."}

    items = []
    ignored_dirs = {".git", "node_modules", "target", "dist", ".venv", "__pycache__", "build", ".idea", ".vscode"}

    for root, dirs, files in os.walk(target_dir):
        dirs[:] = [d for d in dirs if d not in ignored_dirs]
        for f in files:
            if f.startswith(".git"):
                continue
            full_file_path = os.path.join(root, f)
            rel_file_path = os.path.relpath(full_file_path, workspace_root).replace("\\", "/")
            items.append(rel_file_path)
        if len(items) > 100:
            break

    return {
        "path": subpath or "./",
        "file_count": len(items),
        "files": items[:100]
    }

@tool_registry.register(
    name="agent_search_web",
    description="Performs real-time web search for technical specifications, benchmarks, libraries, research papers, or documentation."
)
async def agent_search_web(query: str, max_results: int = 4, context: ToolExecutionContext = None) -> dict[str, Any]:
    results = await searxng_client.search_web(query.strip(), limit=min(max_results, 4))
    for r in results:
        if "snippet" in r:
            r["snippet"] = r["snippet"][:250]
    return {"query": query, "results": results}

@tool_registry.register(
    name="agent_read_link",
    description=(
        "Fetches and extracts full article content, documentation pages, research papers, or YouTube video transcripts.\n"
        "MANDATORY: Call this on any user-provided URLs in Turn 1 before proceeding."
    )
)
async def agent_read_link(url: str, context: ToolExecutionContext = None) -> dict[str, Any]:
    logger.info(f"[agent_read_link] Fetching URL: {url}")

    if is_youtube_url(url):
        canonical_yt_url = normalize_youtube_url(url)
        client, key_idx, active_model = client_manager.get_client_for_model("gemini-3.5-flash-lite")
        if client:
            try:
                prompt_instruction = (
                    "Analyze this YouTube video thoroughly. Provide the exact video title, channel author, and a structured, "
                    "comprehensive technical summary with key takeaways, timestamps for major topic transitions, and findings."
                )
                response = await client.aio.models.generate_content(
                    model=active_model,
                    contents=[
                        types.Part.from_uri(file_uri=canonical_yt_url, mime_type="video/mp4"),
                        prompt_instruction
                    ]
                )
                if response.text:
                    return {
                        "url": canonical_yt_url,
                        "type": "youtube_video_summary",
                        "content": response.text.strip()[:MAX_TOOL_OUTPUT_CHARS]
                    }
            except Exception as e:
                client_manager.report_error(key_idx, active_model, e)
                logger.warning(f"[agent_read_link] YouTube video parse failed: {e}")

    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
        }
        async with httpx.AsyncClient(timeout=10.0, follow_redirects=True) as client:
            resp = await client.get(url, headers=headers)
            if resp.status_code != 200:
                return {"url": url, "error": f"HTTP status {resp.status_code}"}

            extracted_text = trafilatura.extract(
                resp.text,
                include_links=True,
                include_tables=True,
                output_format="markdown"
            )

            if not extracted_text or len(extracted_text.strip()) < 100:
                soup = BeautifulSoup(resp.text, "html.parser")
                for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
                    tag.decompose()
                extracted_text = soup.get_text(separator="\n", strip=True)

            cleaned_text = extracted_text[:MAX_TOOL_OUTPUT_CHARS].strip() if extracted_text else "Unable to extract main content."

            return {
                "url": url,
                "length": len(cleaned_text),
                "content": cleaned_text
            }

    except Exception as e:
        logger.error(f"[agent_read_link] Error reading {url}: {e}")
        return {"url": url, "error": f"Failed to load URL: {str(e)}"}

@tool_registry.register(
    name="agent_search_discord_history",
    description=(
        "Searches recent message history in a Discord channel for user discussions or project lore (uses standard channel permissions).\n"
        "- query: Keyword or topic to search for.\n"
        "- channel_id: Optional channel ID (defaults to current channel).\n"
        "- limit: Number of messages to inspect (max 50)."
    )
)
async def agent_search_discord_history(query: str, channel_id: str = "", limit: int = 30, context: ToolExecutionContext = None) -> dict[str, Any]:
    if not context or not context.bot:
        return {"error": "Discord bot context unavailable."}

    target_channel = context.channel
    if channel_id and channel_id.strip().isdigit():
        try:
            c_id = int(channel_id.strip())
            target_channel = context.bot.get_channel(c_id) or await context.bot.fetch_channel(c_id)
        except Exception:
            pass

    if not target_channel or not isinstance(target_channel, (discord.TextChannel, discord.Thread)):
        return {"error": "Channel does not support message history inspection."}

    matched = []
    q_lower = query.lower().strip()
    try:
        async for m in target_channel.history(limit=min(limit, 50)):
            if not q_lower or q_lower in m.clean_content.lower():
                matched.append({
                    "id": str(m.id),
                    "author": m.author.display_name,
                    "content": m.clean_content[:300],
                    "timestamp": m.created_at.isoformat()
                })
        return {
            "query": query,
            "channel": target_channel.name,
            "matches_found": len(matched),
            "results": matched[:10]
        }
    except Exception as e:
        return {"error": f"Discord history search failed: {e}"}